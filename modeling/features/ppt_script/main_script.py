import os
import sys
import json
from pathlib import Path
from pptx import Presentation
from dotenv import load_dotenv

# =========================================================
# 경로 설정: 직접 실행과 모듈 실행 모두 지원
# =========================================================
current_dir = Path(__file__).resolve().parent
project_root = current_dir.parent.parent

# sys.path에 프로젝트 루트 추가
if str(project_root) not in sys.path:
    sys.path.insert(0, str(project_root))

# =========================================================
# script_llm Import: 직접 실행 시 절대 임포트
# =========================================================
try:
    # 모듈로 실행될 때 (python -m features.ppt_script.main_script)
    from .script_llm import generate_script_and_qna
except ImportError:
    # 직접 실행될 때 (python main_script.py)
    from script_llm import generate_script_and_qna

load_dotenv()

def extract_text_from_pptx(pptx_path: str) -> str:
    """
    PPT 파일에서 텍스트 추출
    
    Args:
        pptx_path: PPT 파일 경로
    
    Returns:
        str: 슬라이드별로 구조화된 텍스트
    """
    if not os.path.exists(pptx_path):
        return None
    
    try:
        prs = Presentation(pptx_path)
        slides_text = []
        
        for i, slide in enumerate(prs.slides):
            # 슬라이드 제목 추출
            title = slide.shapes.title.text.strip() if slide.shapes.title else "무제"
            
            # 슬라이드 내용 추출
            content = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    # 제목은 이미 추출했으므로 제외
                    if shape == slide.shapes.title:
                        continue
                    
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            content.append(text)
            
            # 슬라이드 정보를 구조화된 형태로 저장
            slide_info = f"[[Slide {i+1}]] Title: {title}\nContent: {' '.join(content)}"
            slides_text.append(slide_info)
        
        return "\n\n".join(slides_text)
    except Exception as e:
        print(f"[오류] PPT 파일 읽기 실패: {e}")
        import traceback
        traceback.print_exc()
        return None


def main(pptx_path: str = None):
    """
    Step 4: PPT 발표 대본 및 Q&A 생성 메인 함수
    
    Args:
        pptx_path: PPT 파일 경로 (선택적, 없으면 기본 경로 사용)
    
    Returns:
        dict: 생성된 스크립트 데이터 또는 None
    """
    print("="*60)
    print("[Step 4] PPT 발표 대본 및 Q&A 생성")
    print("="*60)

    # 1. 입력 파일 경로 설정
    if pptx_path:
        # 파라미터로 경로가 주어진 경우
        target_file = Path(pptx_path)
    else:
        # 기본 경로 사용
        input_folder = project_root / "data" / "script_input"
        input_folder.mkdir(parents=True, exist_ok=True)
        
        target_file = input_folder / "ppt_ex.pptx"
        if not target_file.exists():
            pptx_files = list(input_folder.glob("*.pptx"))
            if pptx_files:
                target_file = pptx_files[0]
            else:
                print(f"[!] 오류: '{input_folder}' 폴더에 PPT 파일이 없습니다.")
                print(f"    경로: {input_folder.absolute()}")
                return None
    
    # 파일 존재 확인
    if not target_file.exists():
        print(f"[!] 파일이 존재하지 않습니다: {target_file}")
        return None
    
    print(f"[*] PPT 파일 읽는 중: {target_file.name}")

    # 2. PPT 텍스트 추출
    ppt_text = extract_text_from_pptx(str(target_file))
    
    if not ppt_text: 
        print(f"[!] PPT 파일을 읽을 수 없습니다: {target_file}")
        return None
    
    print(f"[*] PPT 텍스트 추출 완료 (길이: {len(ppt_text)}자)")

    # 3. Gemini 호출 - 대본 및 Q&A 생성
    print("[*] AI 대본 생성 중...")
    json_data = generate_script_and_qna(ppt_text)
    
    # 4. 결과 저장
    if json_data:
        output_folder = project_root / "data" / "report"
        output_folder.mkdir(parents=True, exist_ok=True)
        output_path = output_folder / "script_flow.json"
        
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(json_data, f, ensure_ascii=False, indent=2)
        
        # 결과 요약 출력
        slide_count = len(json_data.get("slides", []))
        qna_count = len(json_data.get("qna", []))
        
        print(f"\n[*] 실전 발표용 데이터 생성 완료!")
        print(f"    - 슬라이드 대본: {slide_count}개")
        print(f"    - 예상 Q&A: {qna_count}개")
        print(f"    - 파일 위치: {output_path}")
        
        print("="*60)
        print("[Step 4 완료]")
        
        return json_data
    else:
        print("[!] 대본 생성 실패")
        print("="*60)
        return None

if __name__ == "__main__":
    main()
